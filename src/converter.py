"""
Document conversion pipeline using Unstructured.io for all format support.
Integrates OCR for scanned content and provides unified output.
Supports MBOX, Google Takeout, LaTeX, and 200+ other formats.
"""

import os
import re
import json
import pandas as pd
from typing import Optional, List, Dict, Any, Tuple
from pathlib import Path
from loguru import logger

# Document processing libraries - unstructured.io
try:
    from unstructured.partition.auto import partition
    from unstructured.partition.pdf import partition_pdf
    from unstructured.partition.image import partition_image
    from unstructured.partition.docx import partition_docx
    from unstructured.partition.pptx import partition_pptx
    from unstructured.partition.html import partition_html
    from unstructured.partition.text import partition_text
    from unstructured.partition.epub import partition_epub
    UNSTRUCTURED_AVAILABLE = True
except ImportError as e:
    logger.warning(f"Unstructured.io not available: {e}. Some features will be limited.")
    UNSTRUCTURED_AVAILABLE = False

from .config import get_config, get_project_root
from .utils import (
    detect_file_type, is_scanned_pdf, get_output_path,
    ensure_dir, ConversionState, write_failed_files, write_summary
)
from .ocr import OCRProcessor
from .cleaner import ContentOptimizer
from .monitor import ProgressTracker

# Knowledge system imports
try:
    from .cache_manager import CacheManager
    from .metadata_extractor import MetadataExtractor
    from .timeline_builder import TimelineBuilder
    from .embedding_engine import EmbeddingEngine
    from .backlink_engine import BacklinkEngine
    from .graph_builder import GraphBuilder
    KNOWLEDGE_SYSTEM_AVAILABLE = True
except ImportError as e:
    logger.warning(f"Knowledge system not available: {e}")
    KNOWLEDGE_SYSTEM_AVAILABLE = False

# Maximum total extracted size from archives (500 MB)
MAX_ARCHIVE_EXTRACT_SIZE = 500 * 1024 * 1024
# Maximum number of entries in an archive
MAX_ARCHIVE_ENTRIES = 10000


def df_to_markdown(df) -> str:
    """Convert pandas DataFrame to markdown table."""
    lines = []

    # Header
    header = " | ".join(str(col) for col in df.columns)
    lines.append(f"| {header} |")
    lines.append("|" + "|".join(["---"] * len(df.columns)) + "|")

    # Data rows
    for _, row in df.iterrows():
        row_str = " | ".join(str(val) if pd.notna(val) else "" for val in row)
        lines.append(f"| {row_str} |")

    return "\n".join(lines)


def _is_safe_archive_member(member_name: str, target_dir: str) -> bool:
    """Check if an archive member path is safe (no path traversal)."""
    # Resolve the full target path
    target = Path(target_dir).resolve()
    member_path = (target / member_name).resolve()
    # Ensure the member path is within the target directory
    try:
        member_path.relative_to(target)
        return True
    except ValueError:
        return False


class DocumentConverterPipeline:
    """
    Main document conversion pipeline.
    Handles all file types and routes to appropriate converter.
    """

    def __init__(self):
        self.config = get_config()
        self.ocr_processor = OCRProcessor()
        self.content_optimizer = ContentOptimizer()
        self.state = ConversionState()

        # Initialize knowledge system components
        if KNOWLEDGE_SYSTEM_AVAILABLE and self.config.knowledge.enabled:
            self.cache_manager = CacheManager()
            self.metadata_extractor = MetadataExtractor()
            self.timeline_builder = TimelineBuilder(self.cache_manager)
            self.embedding_engine = EmbeddingEngine(self.cache_manager)
            self.backlink_engine = BacklinkEngine(self.cache_manager, self.embedding_engine)
            self.graph_builder = GraphBuilder(self.cache_manager)

            # Initialize embedding engine
            self.embedding_engine.initialize()

            logger.info("Knowledge system initialized")
        else:
            self.cache_manager = None
            self.metadata_extractor = None
            self.timeline_builder = None
            self.embedding_engine = None
            self.backlink_engine = None
            self.graph_builder = None

        # Track files processed in current batch (for deferred timeline rebuild)
        self._batch_files_processed: List[str] = []

    def _convert_with_unstructured(self, file_path: str, output_path: str) -> bool:
        """Convert file using Unstructured.io (supports all formats)."""
        if not UNSTRUCTURED_AVAILABLE:
            logger.warning("Unstructured.io not available")
            return False

        try:
            ext = Path(file_path).suffix.lower()
            config = get_config()

            # Use appropriate partition function based on file type
            if ext == '.pdf':
                elements = partition_pdf(file_path, infer_table_structure=True)
            elif ext in ['.docx', '.doc', '.docm']:
                elements = partition_docx(file_path, infer_table_structure=True)
            elif ext in ['.pptx', '.ppt', '.pptm']:
                elements = partition_pptx(file_path, infer_table_structure=True)
            elif ext in ['.html', '.htm', '.xhtml']:
                elements = partition_html(file_path)
            elif ext in ['.txt', '.md', '.markdown', '.rst', '.adoc']:
                elements = partition_text(file_path)
            elif ext == '.epub':
                elements = partition_epub(file_path)
            elif ext in ['.xlsx', '.xls', '.xlsm']:
                elements = partition(file_path, infer_table_structure=True)
            elif ext in config.supported_formats.images:
                elements = partition_image(file_path)
            elif ext in config.supported_formats.code:
                elements = partition_text(file_path)
            elif ext in config.supported_formats.data:
                elements = partition(file_path, infer_table_structure=True)
            else:
                elements = partition(file_path, infer_table_structure=True)

            # Extract text from elements
            text_parts = []
            for element in elements:
                text = element.text.strip()
                if text:
                    text_parts.append(text)

            text = "\n\n".join(text_parts)
            optimized = self.content_optimizer.process(text)

            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(optimized)

            logger.info(f"Successfully converted with Unstructured.io: {file_path}")
            return True

        except Exception as e:
            logger.error(f"Unstructured.io conversion failed: {e}")
            return False

    def convert_file(self, file_path: str) -> Tuple[str, bool]:
        """
        Convert a single file to markdown/text with knowledge system integration.

        Args:
            file_path: Path to input file

        Returns:
            Tuple of (output_path, success)
        """
        file_type = detect_file_type(file_path)

        # Generate output path
        output_path = get_output_path(
            file_path,
            f".{self.config.processing.output_format}"
        )

        # Ensure output directory exists
        ensure_dir(os.path.dirname(output_path))

        try:
            if file_type == "document":
                success = self._convert_document(file_path, output_path)
            elif file_type == "image":
                success = self._convert_image(file_path, output_path)
            elif file_type == "text":
                success = self._convert_text(file_path, output_path)
            elif file_type == "code":
                success = self._convert_code(file_path, output_path)
            elif file_type == "archive":
                success = self._convert_archive(file_path, output_path)
            elif file_type == "audio":
                success = self._convert_audio(file_path, output_path)
            elif file_type == "video":
                success = self._convert_video(file_path, output_path)
            elif file_type == "data":
                success = self._convert_data(file_path, output_path)
            elif file_type == "ebook":
                success = self._convert_ebook(file_path, output_path)
            elif file_type == "spreadsheet":
                success = self._convert_spreadsheet(file_path, output_path)
            elif file_type == "presentation":
                success = self._convert_presentation(file_path, output_path)
            elif file_type == "email":
                success = self._convert_mbox(file_path, output_path)
            elif file_type == "google_takeout":
                success = self._convert_google_takeout(file_path, output_path)
            elif file_type == "latex":
                success = self._convert_latex(file_path, output_path)
            else:
                success = self._convert_with_unstructured(file_path, output_path)

            # Apply knowledge system processing if enabled (but defer timeline rebuild)
            if success and self.cache_manager:
                self._apply_knowledge_system(file_path, output_path, defer_timeline=True)

            # Update state
            if success:
                self.state.mark_completed(file_path)
                self._batch_files_processed.append(file_path)
            else:
                self.state.mark_failed(file_path)

            return output_path, success

        except Exception as e:
            logger.error(f"Conversion failed for {file_path}: {e}")
            self.state.mark_failed(file_path)
            return output_path, False

    def finalize_batch(self):
        """
        Called after a batch of files is processed.
        Rebuilds timeline once instead of per-file.
        """
        if self.timeline_builder and self._batch_files_processed:
            try:
                self.timeline_builder.build_timeline()
                logger.info(f"Timeline rebuilt after batch of {len(self._batch_files_processed)} files")
            except Exception as e:
                logger.error(f"Failed to rebuild timeline after batch: {e}")
        self._batch_files_processed.clear()

    def _apply_knowledge_system(self, file_path: str, output_path: str, defer_timeline: bool = False):
        """
        Apply knowledge system processing to converted file.
        Extracts metadata, generates embeddings, creates backlinks, and enriches markdown.

        Args:
            file_path: Original input file path
            output_path: Converted output file path
            defer_timeline: If True, skip timeline rebuild (caller will do it in batch)
        """
        try:
            # Read converted content
            with open(output_path, 'r', encoding='utf-8') as f:
                content = f.read()

            # Extract metadata
            metadata = self.metadata_extractor.extract_all_metadata(content, file_path)

            # Register file in cache
            file_id = self.cache_manager.register_file(file_path, metadata['title'], content)

            # Store metadata in cache
            date_info_list = [
                {'date_str': d['date_str'], 'date_normalized': d['date_normalized'], 'date_type': d['date_type']}
                for d in metadata['dates']
            ]
            self.cache_manager.store_dates(file_id, date_info_list)
            self.cache_manager.store_people(file_id, metadata['people'])
            self.cache_manager.store_tags(file_id, metadata['tags'])

            # Generate embedding
            self.embedding_engine.add_embedding(file_id, content)

            # Generate backlinks
            self.backlink_engine.generate_backlinks_for_file(file_path, content)

            # Only rebuild timeline immediately if not deferring
            if not defer_timeline:
                self.timeline_builder.build_timeline()

            # Enrich markdown with frontmatter and knowledge sections
            enriched_content = self._enrich_markdown(
                content,
                metadata,
                file_path
            )

            # Write enriched content back
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(enriched_content)

            logger.info(f"Knowledge system processing complete for {file_path}")

        except Exception as e:
            logger.error(f"Knowledge system processing failed for {file_path}: {e}")

    def _enrich_markdown(self, content: str, metadata: Dict[str, Any], file_path: str) -> str:
        """Enrich markdown with frontmatter and knowledge sections."""
        from datetime import datetime

        # Build frontmatter
        frontmatter_lines = ["---"]
        frontmatter_lines.append(f"title: {metadata['title']}")
        frontmatter_lines.append(f"created: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")

        if metadata['tags']:
            frontmatter_lines.append(f"tags: {metadata['tags']}")

        if metadata['people']:
            frontmatter_lines.append(f"people: {metadata['people']}")

        frontmatter_lines.append("---")
        frontmatter_lines.append("")

        # Build knowledge sections
        knowledge_sections = []

        # Timeline section
        timeline_md = self.timeline_builder.generate_timeline_markdown(file_path)
        if timeline_md:
            knowledge_sections.append(timeline_md)

        # Related notes section
        backlinks_md = self.backlink_engine.generate_backlinks_markdown(file_path)
        if backlinks_md:
            knowledge_sections.append(backlinks_md)

        # Mentioned people section
        if metadata['people']:
            people_lines = ["## Mentioned People"]
            for person in metadata['people']:
                people_lines.append(f"- {person}")
            knowledge_sections.append("\n".join(people_lines))

        # Combine all sections
        enriched = "\n".join(frontmatter_lines) + content

        if knowledge_sections:
            enriched += "\n\n" + "\n\n".join(knowledge_sections)

        return enriched

    # ---- Document converters ----

    def _convert_document(self, file_path: str, output_path: str) -> bool:
        """Convert document files (PDF, DOCX, PPTX, EPUB)."""
        ext = Path(file_path).suffix.lower()

        if UNSTRUCTURED_AVAILABLE:
            logger.info(f"Converting with Unstructured.io: {file_path}")
            success = self._convert_with_unstructured(file_path, output_path)
            if success:
                return True
            logger.warning("Unstructured.io failed, trying fallback")

        logger.info(f"Using fallback converter for {file_path}")
        return self._convert_document_fallback(file_path, output_path, ext)

    def _convert_image(self, file_path: str, output_path: str) -> bool:
        """Convert image files using OCR."""
        if UNSTRUCTURED_AVAILABLE:
            success = self._convert_with_unstructured(file_path, output_path)
            if success:
                return True

        try:
            text, success = self.ocr_processor.process_file(file_path)
            if success and text.strip():
                optimized = self.content_optimizer.process(text)
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(optimized)
                return True
        except Exception as e:
            logger.error(f"Image conversion failed: {e}")
        return False

    def _convert_text(self, file_path: str, output_path: str) -> bool:
        """Convert text-based files (TXT, MD, HTML, etc)."""
        ext = Path(file_path).suffix.lower()

        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()

            if ext in ['.html', '.htm', '.xml']:
                content = self._html_to_markdown(content)
            elif ext in ['.json', '.csv']:
                content = self._structured_to_markdown(content, ext)

            optimized = self.content_optimizer.process(content)

            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(optimized)

            return True

        except Exception as e:
            logger.error(f"Text conversion failed: {e}")
            return False

    def _convert_code(self, file_path: str, output_path: str) -> bool:
        """Convert code files to markdown with syntax highlighting."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                code = f.read()

            ext = Path(file_path).suffix.lower().lstrip('.')
            markdown = f"```{ext}\n{code}\n```"
            optimized = self.content_optimizer.process(markdown)

            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(optimized)

            return True

        except Exception as e:
            logger.error(f"Code conversion failed: {e}")
            return False

    def _convert_data(self, file_path: str, output_path: str) -> bool:
        """Convert data files (JSON, YAML, etc.)."""
        if UNSTRUCTURED_AVAILABLE:
            success = self._convert_with_unstructured(file_path, output_path)
            if success:
                return True

        return self._convert_text(file_path, output_path)

    def _convert_ebook(self, file_path: str, output_path: str) -> bool:
        """Convert ebook files (EPUB, MOBI)."""
        if UNSTRUCTURED_AVAILABLE:
            success = self._convert_with_unstructured(file_path, output_path)
            if success:
                return True

        return self._convert_document_fallback(file_path, output_path, Path(file_path).suffix.lower())

    def _convert_spreadsheet(self, file_path: str, output_path: str) -> bool:
        """Convert spreadsheet files (Excel, CSV, TSV)."""
        try:
            ext = Path(file_path).suffix.lower()

            if ext in ['.xlsx', '.xls', '.xlsm']:
                df = pd.read_excel(file_path, sheet_name=None)
            elif ext == '.csv':
                df = pd.read_csv(file_path)
            elif ext == '.tsv':
                df = pd.read_csv(file_path, sep='\t')
            else:
                return False

            if isinstance(df, dict):
                markdown_parts = []
                for sheet_name, sheet_df in df.items():
                    markdown_parts.append(f"## Sheet: {sheet_name}\n")
                    markdown_parts.append(df_to_markdown(sheet_df))
                text = "\n\n".join(markdown_parts)
            else:
                text = df_to_markdown(df)

            optimized = self.content_optimizer.process(text)

            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(optimized)

            return True

        except Exception as e:
            logger.error(f"Spreadsheet conversion failed: {e}")
            return False

    def _convert_presentation(self, file_path: str, output_path: str) -> bool:
        """Convert presentation files (PPTX, PPT)."""
        if UNSTRUCTURED_AVAILABLE:
            success = self._convert_with_unstructured(file_path, output_path)
            if success:
                return True

        return self._convert_document_fallback(file_path, output_path, Path(file_path).suffix.lower())

    def _convert_archive(self, file_path: str, output_path: str) -> bool:
        """
        Extract and convert archive contents.
        Includes path traversal protection and size limits.
        """
        try:
            import zipfile
            import tarfile
            import shutil

            temp_dir = output_path + "_extracted"
            Path(temp_dir).mkdir(parents=True, exist_ok=True)

            ext = Path(file_path).suffix.lower()
            combined_name = Path(file_path).stem
            # Handle .tar.gz / .tgz
            if file_path.endswith('.tar.gz'):
                ext = '.tar.gz'

            if ext == '.zip':
                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                    # Validate members before extraction
                    total_size = 0
                    for info in zip_ref.infolist():
                        if not _is_safe_archive_member(info.filename, temp_dir):
                            logger.error(f"Unsafe archive member path: {info.filename}")
                            shutil.rmtree(temp_dir, ignore_errors=True)
                            return False
                        total_size += info.file_size
                        if total_size > MAX_ARCHIVE_EXTRACT_SIZE:
                            logger.error(f"Archive exceeds max extract size ({MAX_ARCHIVE_EXTRACT_SIZE} bytes)")
                            shutil.rmtree(temp_dir, ignore_errors=True)
                            return False
                    if len(zip_ref.infolist()) > MAX_ARCHIVE_ENTRIES:
                        logger.error(f"Archive has too many entries (>{MAX_ARCHIVE_ENTRIES})")
                        shutil.rmtree(temp_dir, ignore_errors=True)
                        return False
                    zip_ref.extractall(temp_dir)

            elif ext in ['.tar', '.tar.gz', '.tgz']:
                with tarfile.open(file_path, 'r:*') as tar_ref:
                    members = tar_ref.getmembers()
                    total_size = sum(m.size for m in members if m.isfile())
                    if total_size > MAX_ARCHIVE_EXTRACT_SIZE:
                        logger.error(f"Archive exceeds max extract size ({MAX_ARCHIVE_EXTRACT_SIZE} bytes)")
                        shutil.rmtree(temp_dir, ignore_errors=True)
                        return False
                    if len(members) > MAX_ARCHIVE_ENTRIES:
                        logger.error(f"Archive has too many entries (>{MAX_ARCHIVE_ENTRIES})")
                        shutil.rmtree(temp_dir, ignore_errors=True)
                        return False
                    # Filter out unsafe members
                    safe_members = [
                        m for m in members
                        if _is_safe_archive_member(m.name, temp_dir)
                    ]
                    if len(safe_members) < len(members):
                        logger.warning(f"Skipped {len(members) - len(safe_members)} unsafe archive members")
                    tar_ref.extractall(temp_dir, members=safe_members)

            elif ext == '.7z':
                try:
                    import py7zr
                    with py7zr.SevenZipFile(file_path, mode='r') as z:
                        z.extractall(path=temp_dir)
                except ImportError:
                    logger.warning("7z extraction requires py7zr — install with: pip install py7zr")
                    shutil.rmtree(temp_dir, ignore_errors=True)
                    return False

            elif ext == '.rar':
                try:
                    import rarfile
                    with rarfile.RarFile(file_path) as rf:
                        rf.extractall(temp_dir)
                except ImportError:
                    logger.warning("RAR extraction requires rarfile — install with: pip install rarfile")
                    shutil.rmtree(temp_dir, ignore_errors=True)
                    return False
            else:
                shutil.rmtree(temp_dir, ignore_errors=True)
                return False

            # Recursively convert extracted files
            extracted_files = sorted(Path(temp_dir).rglob('*'))
            converted_texts = []
            for extracted_file in extracted_files:
                if extracted_file.is_file():
                    sub_output, sub_success = self.convert_file(str(extracted_file))
                    if sub_success and os.path.exists(sub_output):
                        with open(sub_output, 'r', encoding='utf-8', errors='ignore') as f:
                            converted_texts.append(f"## {extracted_file.name}\n\n{f.read()}")

            # Write combined output
            if converted_texts:
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(f"# Archive: {combined_name}\n\n" + "\n\n---\n\n".join(converted_texts))

            # Clean up temp directory
            shutil.rmtree(temp_dir, ignore_errors=True)

            logger.info(f"Extracted and converted archive: {file_path}")
            return bool(converted_texts)

        except Exception as e:
            logger.error(f"Archive extraction failed: {e}")
            return False

    def _convert_audio(self, file_path: str, output_path: str) -> bool:
        """Convert audio file using Whisper for transcription."""
        try:
            import whisper
            model = whisper.load_model("base")
            result = model.transcribe(file_path)
            text = result.get("text", "")
            if text.strip():
                optimized = self.content_optimizer.process(text)
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(optimized)
                return True
            return False
        except ImportError:
            logger.warning("Audio transcription requires openai-whisper — install with: pip install openai-whisper")
            return False
        except Exception as e:
            logger.error(f"Audio conversion failed: {e}")
            return False

    def _convert_video(self, file_path: str, output_path: str) -> bool:
        """Convert video file — extract audio and transcribe."""
        try:
            import subprocess
            import tempfile

            # Extract audio track using ffmpeg
            with tempfile.NamedTemporaryFile(suffix='.wav', delete=False) as tmp:
                tmp_audio = tmp.name

            result = subprocess.run(
                ['ffmpeg', '-i', file_path, '-vn', '-acodec', 'pcm_s16le',
                 '-ar', '16000', '-ac', '1', tmp_audio, '-y'],
                capture_output=True, text=True, timeout=300
            )

            if result.returncode == 0 and os.path.exists(tmp_audio):
                success = self._convert_audio(tmp_audio, output_path)
                os.unlink(tmp_audio)
                return success

            os.unlink(tmp_audio)
            logger.warning(f"ffmpeg audio extraction failed for {file_path}")
            return False

        except FileNotFoundError:
            logger.warning("Video processing requires ffmpeg to be installed")
            return False
        except Exception as e:
            logger.error(f"Video conversion failed: {e}")
            return False

    # ---- New format converters ----

    def _convert_mbox(self, file_path: str, output_path: str) -> bool:
        """Convert MBOX email archive to markdown."""
        try:
            import mailbox
            from email import policy
            from email.utils import parsedate_to_datetime

            mbox = mailbox.mbox(file_path)
            emails_md = []
            count = 0

            for message in mbox:
                count += 1
                try:
                    subject = message.get('Subject', '(No Subject)')
                    sender = message.get('From', 'Unknown')
                    to = message.get('To', '')
                    date_str = message.get('Date', '')
                    message_id = message.get('Message-ID', '')

                    # Parse date
                    date_display = date_str
                    try:
                        if date_str:
                            dt = parsedate_to_datetime(date_str)
                            date_display = dt.strftime('%Y-%m-%d %H:%M:%S')
                    except Exception:
                        pass

                    # Extract body
                    body = self._extract_email_body(message)

                    email_md = f"## Email: {subject}\n\n"
                    email_md += f"- **From:** {sender}\n"
                    email_md += f"- **To:** {to}\n"
                    email_md += f"- **Date:** {date_display}\n"
                    if message_id:
                        email_md += f"- **Message-ID:** {message_id}\n"
                    email_md += f"\n{body}\n"

                    emails_md.append(email_md)

                except Exception as e:
                    logger.debug(f"Failed to parse email #{count}: {e}")
                    continue

            mbox.close()

            if emails_md:
                header = f"# MBOX Archive: {Path(file_path).stem}\n\n"
                header += f"Total emails: {len(emails_md)}\n\n---\n\n"
                full_text = header + "\n---\n\n".join(emails_md)
                optimized = self.content_optimizer.process(full_text)

                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(optimized)

                logger.info(f"Converted {len(emails_md)} emails from MBOX: {file_path}")
                return True

            return False

        except Exception as e:
            logger.error(f"MBOX conversion failed: {e}")
            return False

    def _extract_email_body(self, message) -> str:
        """Extract text body from an email message."""
        body_parts = []

        if message.is_multipart():
            for part in message.walk():
                content_type = part.get_content_type()
                content_disposition = str(part.get("Content-Disposition", ""))

                if "attachment" in content_disposition:
                    continue

                if content_type == 'text/plain':
                    try:
                        payload = part.get_payload(decode=True)
                        if payload:
                            charset = part.get_content_charset() or 'utf-8'
                            body_parts.append(payload.decode(charset, errors='ignore'))
                    except Exception:
                        pass
                elif content_type == 'text/html':
                    try:
                        payload = part.get_payload(decode=True)
                        if payload:
                            charset = part.get_content_charset() or 'utf-8'
                            html_text = payload.decode(charset, errors='ignore')
                            body_parts.append(self._html_to_markdown(html_text))
                    except Exception:
                        pass
        else:
            content_type = message.get_content_type()
            try:
                payload = message.get_payload(decode=True)
                if payload:
                    charset = message.get_content_charset() or 'utf-8'
                    text = payload.decode(charset, errors='ignore')
                    if content_type == 'text/html':
                        text = self._html_to_markdown(text)
                    body_parts.append(text)
            except Exception:
                pass

        return "\n\n".join(body_parts) if body_parts else "(No body content)"

    def _convert_google_takeout(self, file_path: str, output_path: str) -> bool:
        """
        Convert Google Takeout JSON files to markdown.
        Handles common Takeout structures: activity, contacts, chat messages, etc.
        """
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                data = json.load(f)

            markdown_parts = [f"# Google Takeout: {Path(file_path).stem}\n"]

            if isinstance(data, list):
                for i, item in enumerate(data):
                    md = self._takeout_item_to_markdown(item, i + 1)
                    if md:
                        markdown_parts.append(md)
            elif isinstance(data, dict):
                # Some Takeout files have a top-level key with the list
                for key, value in data.items():
                    if isinstance(value, list):
                        markdown_parts.append(f"\n## {key}\n")
                        for i, item in enumerate(value):
                            md = self._takeout_item_to_markdown(item, i + 1)
                            if md:
                                markdown_parts.append(md)
                    elif isinstance(value, dict):
                        markdown_parts.append(f"\n## {key}\n")
                        markdown_parts.append(self._dict_to_markdown_list(value))
                    else:
                        markdown_parts.append(f"- **{key}:** {value}")

            full_text = "\n\n".join(markdown_parts)
            optimized = self.content_optimizer.process(full_text)

            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(optimized)

            logger.info(f"Converted Google Takeout file: {file_path}")
            return True

        except json.JSONDecodeError:
            logger.debug(f"Not valid JSON, trying as regular text: {file_path}")
            return self._convert_text(file_path, output_path)
        except Exception as e:
            logger.error(f"Google Takeout conversion failed: {e}")
            return False

    def _takeout_item_to_markdown(self, item: Any, index: int) -> str:
        """Convert a single Google Takeout item to markdown."""
        if isinstance(item, dict):
            lines = [f"### Entry {index}"]
            for key, value in item.items():
                if isinstance(value, (list, dict)):
                    lines.append(f"- **{key}:** {json.dumps(value, ensure_ascii=False, default=str)[:500]}")
                else:
                    lines.append(f"- **{key}:** {value}")
            return "\n".join(lines)
        elif isinstance(item, str):
            return f"- {item}"
        return ""

    def _dict_to_markdown_list(self, d: Dict[str, Any]) -> str:
        """Convert a dict to a markdown bullet list."""
        lines = []
        for key, value in d.items():
            if isinstance(value, (list, dict)):
                lines.append(f"- **{key}:** {json.dumps(value, ensure_ascii=False, default=str)[:500]}")
            else:
                lines.append(f"- **{key}:** {value}")
        return "\n".join(lines)

    def _convert_latex(self, file_path: str, output_path: str) -> bool:
        """Convert LaTeX files to markdown."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()

            # Strip LaTeX commands to produce readable text
            md = self._latex_to_markdown(content)
            optimized = self.content_optimizer.process(md)

            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(optimized)

            logger.info(f"Converted LaTeX file: {file_path}")
            return True

        except Exception as e:
            logger.error(f"LaTeX conversion failed: {e}")
            return False

    def _latex_to_markdown(self, latex: str) -> str:
        """Basic LaTeX to markdown conversion."""
        text = latex

        # Remove comments
        text = re.sub(r'(?<!\\)%.*$', '', text, flags=re.MULTILINE)

        # Convert sections
        text = re.sub(r'\\section\{([^}]*)\}', r'# \1', text)
        text = re.sub(r'\\subsection\{([^}]*)\}', r'## \1', text)
        text = re.sub(r'\\subsubsection\{([^}]*)\}', r'### \1', text)

        # Convert formatting
        text = re.sub(r'\\textbf\{([^}]*)\}', r'**\1**', text)
        text = re.sub(r'\\textit\{([^}]*)\}', r'*\1*', text)
        text = re.sub(r'\\emph\{([^}]*)\}', r'*\1*', text)
        text = re.sub(r'\\underline\{([^}]*)\}', r'\1', text)
        text = re.sub(r'\\texttt\{([^}]*)\}', r'`\1`', text)

        # Convert lists
        text = re.sub(r'\\begin\{itemize\}', '', text)
        text = re.sub(r'\\end\{itemize\}', '', text)
        text = re.sub(r'\\begin\{enumerate\}', '', text)
        text = re.sub(r'\\end\{enumerate\}', '', text)
        text = re.sub(r'\\item\s*', '- ', text)

        # Convert math (keep inline math as-is with $ delimiters)
        text = re.sub(r'\\\[', '$$', text)
        text = re.sub(r'\\\]', '$$', text)

        # Remove common preamble commands
        text = re.sub(r'\\documentclass\{[^}]*\}', '', text)
        text = re.sub(r'\\usepackage(\[[^\]]*\])?\{[^}]*\}', '', text)
        text = re.sub(r'\\begin\{document\}', '', text)
        text = re.sub(r'\\end\{document\}', '', text)
        text = re.sub(r'\\title\{([^}]*)\}', r'# \1', text)
        text = re.sub(r'\\author\{([^}]*)\}', r'*Author: \1*', text)
        text = re.sub(r'\\date\{([^}]*)\}', r'*Date: \1*', text)
        text = re.sub(r'\\maketitle', '', text)

        # Remove remaining simple commands
        text = re.sub(r'\\[a-zA-Z]+\{([^}]*)\}', r'\1', text)
        text = re.sub(r'\\[a-zA-Z]+', '', text)

        # Clean up extra whitespace
        text = re.sub(r'\n{3,}', '\n\n', text)

        return text.strip()

    # ---- Fallback converters ----

    def _convert_document_fallback(self, file_path: str, output_path: str, ext: str) -> bool:
        """Fallback document conversion when Unstructured.io is unavailable."""
        config = get_config()

        try:
            if ext in config.supported_formats.archives:
                return self._convert_archive(file_path, output_path)
            elif ext in config.supported_formats.audio:
                return self._convert_audio(file_path, output_path)
            elif ext in config.supported_formats.video:
                return self._convert_video(file_path, output_path)
            elif ext in config.supported_formats.code:
                return self._convert_code(file_path, output_path)
            elif ext in ['.docx', '.doc', '.docm']:
                return self._convert_docx_fallback(file_path, output_path)
            elif ext in ['.pptx', '.ppt', '.pptm']:
                return self._convert_pptx_fallback(file_path, output_path)
            elif ext in config.supported_formats.ebooks:
                return self._convert_epub_fallback(file_path, output_path)
            elif ext == '.pdf':
                return self._convert_pdf_fallback(file_path, output_path)
            elif ext in config.supported_formats.spreadsheets:
                return self._convert_spreadsheet(file_path, output_path)
            elif ext in config.supported_formats.text:
                return self._convert_text(file_path, output_path)
            else:
                logger.warning(f"Unsupported format for fallback: {ext}")
                return False
        except Exception as e:
            logger.error(f"Fallback conversion failed: {e}")
            return False

    def _convert_docx_fallback(self, file_path: str, output_path: str) -> bool:
        """Convert DOCX using python-docx."""
        try:
            from docx import Document

            doc = Document(file_path)
            text_parts = []

            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    text_parts.append(row_text)

            text = "\n\n".join(text_parts)
            optimized = self.content_optimizer.process(text)

            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(optimized)

            return True
        except ImportError:
            logger.error("python-docx not installed")
            return False
        except Exception as e:
            logger.error(f"DOCX conversion failed: {e}")
            return False

    def _convert_pptx_fallback(self, file_path: str, output_path: str) -> bool:
        """Convert PPTX using python-pptx."""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                slide_texts.append(paragraph.text)
                if slide_texts:
                    text_parts.append(f"## Slide {slide_num}\n\n" + "\n\n".join(slide_texts))

            text = "\n\n---\n\n".join(text_parts)
            optimized = self.content_optimizer.process(text)

            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(optimized)

            return True
        except ImportError:
            logger.error("python-pptx not installed")
            return False
        except Exception as e:
            logger.error(f"PPTX conversion failed: {e}")
            return False

    def _convert_epub_fallback(self, file_path: str, output_path: str) -> bool:
        """Convert EPUB using ebooklib."""
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup

            book = epub.read_epub(file_path)
            text_parts = []

            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    text = soup.get_text()
                    if text.strip():
                        text_parts.append(text)

            text = "\n\n".join(text_parts)
            optimized = self.content_optimizer.process(text)

            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(optimized)

            return True
        except ImportError:
            logger.error("ebooklib or beautifulsoup4 not installed")
            return False
        except Exception as e:
            logger.error(f"EPUB conversion failed: {e}")
            return False

    def _convert_pdf_fallback(self, file_path: str, output_path: str) -> bool:
        """Convert PDF using PyMuPDF as fallback."""
        try:
            import fitz  # PyMuPDF

            with fitz.open(file_path) as doc:
                text_parts = []
                for page in doc:
                    text = page.get_text()
                    if text.strip():
                        text_parts.append(text)

            text = "\n\n".join(text_parts)
            optimized = self.content_optimizer.process(text)

            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(optimized)

            return True
        except ImportError:
            logger.error("pymupdf not installed")
            return False
        except Exception as e:
            logger.error(f"PDF fallback conversion failed: {e}")
            return False

    # ---- Helper methods ----

    def _html_to_markdown(self, html_content: str) -> str:
        """Convert HTML content to Markdown."""
        try:
            import html2text
            h = html2text.HTML2Text()
            h.ignore_links = False
            h.ignore_images = True
            h.ignore_tables = False
            h.body_width = 0
            return h.handle(html_content)
        except ImportError:
            from html.parser import HTMLParser

            class MLStripper(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.reset()
                    self.fed: List[str] = []

                def handle_data(self, d):
                    self.fed.append(d)

                def get_data(self):
                    return ''.join(self.fed)

            s = MLStripper()
            s.feed(html_content)
            return s.get_data()

    def _structured_to_markdown(self, content: str, ext: str) -> str:
        """Convert structured data (JSON, CSV) to markdown."""
        try:
            if ext == '.json':
                data = json.loads(content)

                if isinstance(data, list) and len(data) > 0 and isinstance(data[0], dict):
                    headers = list(data[0].keys())
                    rows = [[str(row.get(h, '')) for h in headers] for row in data]

                    lines = []
                    lines.append('| ' + ' | '.join(headers) + ' |')
                    lines.append('|' + '|'.join(['---'] * len(headers)) + '|')
                    for row in rows[:100]:
                        lines.append('| ' + ' | '.join(row) + ' |')

                    return '\n'.join(lines)
                else:
                    return f"```json\n{json.dumps(data, indent=2, ensure_ascii=False, default=str)}\n```"

            elif ext == '.csv':
                import csv
                from io import StringIO

                reader = csv.reader(StringIO(content))
                rows = list(reader)

                if rows:
                    lines = []
                    lines.append('| ' + ' | '.join(rows[0]) + ' |')
                    lines.append('|' + '|'.join(['---'] * len(rows[0])) + '|')
                    for row in rows[1:100]:
                        lines.append('| ' + ' | '.join(row) + ' |')

                    return '\n'.join(lines)

        except Exception as e:
            logger.debug(f"Structured conversion failed: {e}")

        return content


class BatchProcessor:
    """Handles batch processing of files with progress tracking."""

    def __init__(self):
        self.converter = DocumentConverterPipeline()
        self.config = get_config()

    def process_batch(self, files: List[str], tracker: ProgressTracker) -> Tuple[int, int, int]:
        """
        Process a batch of files.

        Returns:
            Tuple of (successful, failed, total_text_bytes)
        """
        successful = 0
        failed = 0
        total_text = 0

        batch_size = self.config.processing.batch_size

        for i, file_path in enumerate(files):
            tracker.start_file(file_path)

            try:
                output_path, success = self.converter.convert_file(file_path)

                if success:
                    successful += 1
                    try:
                        text_size = os.path.getsize(output_path)
                        total_text += text_size
                    except OSError:
                        pass
                else:
                    failed += 1

                tracker.complete_file(success, total_text)

                # Batch delay for thermal management
                if (i + 1) % batch_size == 0 and i < len(files) - 1:
                    import time
                    time.sleep(self.config.processing.batch_delay)
                    # Finalize batch (rebuild timeline once)
                    self.converter.finalize_batch()

            except Exception as e:
                logger.error(f"Unexpected error processing {file_path}: {e}")
                tracker.complete_file(False)
                failed += 1

        # Final batch finalization
        self.converter.finalize_batch()

        return successful, failed, total_text

    def run(self, files: List[str], tracker: ProgressTracker) -> Dict[str, Any]:
        """Run full batch processing pipeline."""
        state = ConversionState()
        pending = state.get_pending_files(files)

        if not pending:
            logger.info("All files already processed")
            return {
                "total": len(files),
                "successful": len(state.completed),
                "failed": len(state.failed),
                "skipped": len(files) - len(state.completed) - len(state.failed)
            }

        logger.info(f"Processing {len(pending)} pending files out of {len(files)} total")

        successful, failed, total_text = self.process_batch(pending, tracker)

        stats = tracker.get_stats()
        write_summary(
            total=len(files),
            successful=stats['successful'],
            failed=stats['failed'],
            total_text=total_text
        )

        if state.failed:
            write_failed_files(list(state.failed))

        return {
            "total": len(files),
            "successful": stats['successful'],
            "failed": stats['failed'],
            "total_text_bytes": total_text
        }


def convert_single_file(file_path: str) -> Tuple[str, bool]:
    """Convert a single file (convenience function)."""
    pipeline = DocumentConverterPipeline()
    return pipeline.convert_file(file_path)


def convert_directory(input_dir: Optional[str] = None) -> Dict[str, Any]:
    """Convert all supported files in directory."""
    from .utils import get_input_files
    from .monitor import ProgressTracker

    files = get_input_files()

    if not files:
        logger.warning("No files found to process")
        return {"total": 0, "successful": 0, "failed": 0}

    tracker = ProgressTracker(len(files))
    processor = BatchProcessor()

    return processor.run(files, tracker)
